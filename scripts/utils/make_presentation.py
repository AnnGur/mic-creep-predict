"""
Generate PowerPoint presentation for the Vivli AMR Surveillance Challenge 2026.
Run: .venv/bin/python scripts/make_presentation.py
Output: reports/mic_creep_predict_presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
OUT = PROJECT_ROOT / "reports" / "mic_creep_predict_presentation.pptx"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x4F, 0x8A)
RED       = RGBColor(0xC0, 0x39, 0x2B)
TEAL      = RGBColor(0x17, 0x7E, 0x89)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x1A, 0x2E)
GRAY      = RGBColor(0x5D, 0x6D, 0x7E)
AMBER     = RGBColor(0xE6, 0x7E, 0x22)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def rect(slide, left, top, width, height, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def textbox(slide, left, top, width, height,
            text, size=18, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_paragraph(tf, text, size=16, bold=False, color=DARK,
                  align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def bullet_box(slide, left, top, width, height,
               items, size=17, color=DARK, title=None, title_color=NAVY):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 1)
        run.font.bold = True
        run.font.color.rgb = title_color
        p.space_before = Pt(0)

    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txb


# ── Slide helpers ─────────────────────────────────────────────────────────────

def header_bar(slide, title, subtitle=None):
    """Navy bar across the top with white title text."""
    rect(slide, 0, 0, 13.33, 1.4, NAVY)
    textbox(slide, 0.4, 0.1, 12.5, 0.75,
            title, size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        textbox(slide, 0.4, 0.88, 12.5, 0.45,
                subtitle, size=16, bold=False, color=RGBColor(0xAA, 0xBB, 0xCC),
                align=PP_ALIGN.LEFT)


def accent_bar(slide, left, top, height=0.06, width=1.6):
    """Small teal accent underline."""
    rect(slide, left, top, width, height, TEAL)


# =============================================================================
# SLIDES
# =============================================================================

def slide_01_title(prs):
    """Title slide."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, DARK)
    rect(sl, 0, 0, 13.33, 0.12, TEAL)
    rect(sl, 0, 7.38, 13.33, 0.12, TEAL)

    rect(sl, 0, 2.7, 13.33, 2.5, RGBColor(0x0D, 0x1B, 0x2A))

    textbox(sl, 0.7, 0.6, 12, 0.7,
            "Vivli AMR Surveillance Data Challenge 2026",
            size=18, color=TEAL, align=PP_ALIGN.LEFT)

    textbox(sl, 0.7, 1.35, 11.8, 1.3,
            "Predicting MIC Creep in\nK. pneumoniae and A. baumannii",
            size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    textbox(sl, 0.7, 2.85, 11.5, 0.55,
            "A machine learning approach to detect pre-resistance MIC drift "
            "before susceptibility thresholds are breached",
            size=17, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.LEFT)

    textbox(sl, 0.7, 3.65, 11.5, 0.45,
            "Antibiotic: Meropenem  |  Dataset: ATLAS (Pfizer/Vivli)  |  2004-2022",
            size=15, color=GRAY, align=PP_ALIGN.LEFT)

    textbox(sl, 0.7, 6.55, 11.5, 0.28,
            "Lead Researcher / Domain Expert: Inna Kucherova  |  Technical Implementation: Anna Gurina",
            size=13, color=GRAY, align=PP_ALIGN.LEFT)
    textbox(sl, 0.7, 6.85, 11.5, 0.28,
            "Proof of Concept submission  |  Open-source code available on GitHub",
            size=12, color=RGBColor(0x4A, 0x5A, 0x6A), align=PP_ALIGN.LEFT)


def slide_02_problem(prs):
    """What is MIC Creep?"""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "The Problem: What is MIC Creep?",
               "A silent precursor to full antibiotic resistance")

    # Left column - definition
    rect(sl, 0.4, 1.6, 5.8, 5.5, WHITE)
    bullet_box(sl, 0.65, 1.75, 5.3, 5.2,
               title="MIC Creep Defined",
               items=[
                   "MIC = Minimum Inhibitory Concentration",
                   "  the lowest drug concentration that stops growth",
                   "",
                   "MIC Creep = gradual upward drift of MIC values",
                   "  across a bacterial population year-over-year",
                   "",
                   "Danger zone: values remain technically 'Susceptible'",
                   "  while drifting toward the Resistant threshold",
                   "",
                   "By the time resistance is officially declared,",
                   "  clinical failures may already be occurring",
               ],
               size=15, color=DARK)

    # Right column - key numbers: K. pneumoniae
    rect(sl, 6.6, 1.6, 6.3, 2.4, NAVY)
    textbox(sl, 6.85, 1.75, 5.8, 0.4,
            "K. pneumoniae + Meropenem (ATLAS)", size=14, bold=True,
            color=WHITE)
    textbox(sl, 6.85, 2.2, 5.8, 1.65,
            "+1.97 mg/L per year\nMIC90 slope (2004-2018)\n\n"
            "Resistance rate: 5% (2007) -> 20% (2024)",
            size=16, bold=False, color=RGBColor(0xBB, 0xDD, 0xFF))

    # Right column - A. baumannii
    rect(sl, 6.6, 4.15, 6.3, 2.0, RED)
    textbox(sl, 6.85, 4.28, 5.8, 0.4,
            "A. baumannii + Meropenem (ATLAS)", size=14, bold=True,
            color=WHITE)
    textbox(sl, 6.85, 4.72, 5.8, 1.3,
            "MIC90 at panel ceiling (32 mg/L) since 2005\n\n"
            "Resistance rate: 39% (2006) -> 69% (2022)\n"
            "Already at crisis level throughout dataset",
            size=15, bold=False, color=RGBColor(0xFF, 0xCC, 0xCC))

    # EUCAST clinical breakpoints table
    eucast_img = PROJECT_ROOT / "EUCAST Clinical Breakpoints.png"
    if eucast_img.exists():
        sl.shapes.add_picture(str(eucast_img), Inches(0.4), Inches(6.22), Inches(12.5), Inches(1.12))
    else:
        rect(sl, 0.4, 6.22, 12.5, 1.12, WHITE)
        textbox(sl, 0.6, 6.65, 12.0, 0.35,
                "EUCAST Clinical Breakpoints - Meropenem (2024): S ≤ 2 mg/L | I > 2-8 mg/L | R > 8 mg/L",
                size=13, color=GRAY, italic=True)


def slide_03_why_matters(prs):
    """Clinical stakes."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Why It Matters",
               "Meropenem is a last-resort carbapenem antibiotic")

    cards = [
        (RED,   "Last Resort Drug",
         "When meropenem fails, treatment options are\n"
         "colistin (nephrotoxic) or ceftazidime-avibactam\n"
         "(increasingly compromised by NDM spread)"),
        (AMBER, "NDM: The Fast-Rising Threat",
         "NDM carbapenemase bypasses avibactam combinations.\n"
         "In training data (2004-2018) NDM is rare.\n"
         "By 2022 it overtakes KPC as the dominant mechanism."),
        (NAVY,  "Combat-Wound Relevance",
         "A. baumannii dominates trauma and wound infections.\n"
         "The military proxy group (wound+male+18-60)\n"
         "represents 2,813 isolates in the ATLAS dataset."),
        (TEAL,  "Surveillance Gap",
         "Standard S/I/R reporting misses MIC creep.\n"
         "An isolate at MIC=4 and MIC=0.25 are both\n"
         "'Susceptible' - yet represent very different risk."),
    ]

    for i, (color, title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 6.5
        y = 1.65 + row * 2.7
        rect(sl, x, y, 6.1, 2.45, color)
        textbox(sl, x + 0.2, y + 0.15, 5.7, 0.5,
                title, size=17, bold=True, color=WHITE)
        textbox(sl, x + 0.2, y + 0.65, 5.7, 1.65,
                body, size=14, color=RGBColor(0xEE, 0xEE, 0xEE))


def slide_04_pathogens(prs):
    """Target pathogens."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Target Pathogens",
               "Both are WHO Critical Priority pathogens for carbapenem resistance")

    # K. pneumoniae
    rect(sl, 0.4, 1.6, 6.0, 5.55, NAVY)
    textbox(sl, 0.65, 1.75, 5.5, 0.55,
            "Klebsiella pneumoniae", size=22, bold=True, color=WHITE)
    textbox(sl, 0.65, 2.35, 5.5, 0.35,
            "Enterobacterales | Gram-negative rod", size=13,
            color=RGBColor(0xAA, 0xCC, 0xFF), italic=True)
    bullet_box(sl, 0.65, 2.8, 5.5, 4.1,
               items=[
                   "Primary cause of carbapenem-resistant nosocomial infections globally",
                   "",
                   "Resistance mechanisms: KPC (dominant 2004-2018), NDM (rising post-2018), OXA-48",
                   "",
                   "EUCAST R breakpoint: MIC > 8 mg/L",
                   "",
                   "89,572 ATLAS isolates (2004-2022)",
               ],
               size=14, color=RGBColor(0xDD, 0xEE, 0xFF))

    # A. baumannii
    rect(sl, 6.9, 1.6, 6.0, 5.55, RGBColor(0x12, 0x5A, 0x50))
    textbox(sl, 7.15, 1.75, 5.5, 0.55,
            "Acinetobacter baumannii", size=22, bold=True, color=WHITE)
    textbox(sl, 7.15, 2.35, 5.5, 0.35,
            "Gammaproteobacteria | Gram-negative coccobacillus", size=13,
            color=RGBColor(0xAA, 0xFF, 0xDD), italic=True)
    bullet_box(sl, 7.15, 2.8, 5.5, 4.1,
               items=[
                   "Primary pathogen in trauma and combat-wound infections",
                   "",
                   "Resistance mechanisms: OXA-23/OXA-40/OXA-58 dominate throughout; NDM secondary",
                   "",
                   "EUCAST R breakpoint: MIC > 8 mg/L",
                   "",
                   "37,540 ATLAS isolates (2004-2022)",
               ],
               size=14, color=RGBColor(0xAA, 0xFF, 0xDD))

    textbox(sl, 0.4, 7.2, 12.5, 0.25,
            "Same modeling pipeline runs for both species - predictions available independently via species parameter in all API endpoints",
            size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


def slide_05_data(prs):
    """Data source."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Data Source: ATLAS",
               "Antimicrobial Testing Leadership and Surveillance - Pfizer/Vivli")

    # Species isolate table
    rect(sl, 0.4, 1.58, 12.5, 1.65, WHITE)

    # Column headers
    col_xs = [0.55, 4.2, 6.9, 9.3, 11.4]
    col_ws = [3.5,  2.5, 2.2, 2.0, 1.7]
    headers = ["Species", "Isolates", "Years", "Countries", "Gene families"]
    for hdr, x, w in zip(headers, col_xs, col_ws):
        textbox(sl, x, 1.63, w, 0.32, hdr, size=12, bold=True,
                color=GRAY, align=PP_ALIGN.LEFT)

    rows = [
        ("K. pneumoniae", "89,572", "2004-2022", "81", "6", NAVY),
        ("A. baumannii",  "37,540", "2004-2022", "79", "6", RED),
    ]
    for ri, (species, isolates, years, countries, genes, color) in enumerate(rows):
        y = 1.98 + ri * 0.62
        rect(sl, 0.4, y, 12.5, 0.58, RGBColor(0xF0, 0xF4, 0xF8) if ri % 2 == 0 else WHITE)
        textbox(sl, col_xs[0], y + 0.1, col_ws[0], 0.38,
                species, size=15, bold=True, color=color, italic=True)
        for val, x, w in zip([isolates, years, countries, genes], col_xs[1:], col_ws[1:]):
            textbox(sl, x, y + 0.1, w, 0.38, val, size=15, bold=True,
                    color=DARK, align=PP_ALIGN.CENTER)

    # Combined total note
    textbox(sl, 0.55, 3.25, 6.0, 0.28,
            "Total: 127,112 isolates across both species",
            size=12, color=GRAY, italic=True)

    rect(sl, 0.4, 3.55, 6.0, 3.6, WHITE)
    bullet_box(sl, 0.6, 3.65, 5.7, 3.4,
               title="What ATLAS Provides",
               items=[
                   "Patient-level MIC measurements",
                   "Country, year, sex, age group",
                   "Specimen source (wound, blood, urine, respiratory)",
                   "Carbapenemase gene results: KPC, NDM, OXA, VIM, IMP, GES",
                   "Hospital-acquired vs community-acquired flag",
                   "",
                   "Access: Vivli AMR Register (controlled, DUA required)",
                   "Raw data: local only, never committed to Git",
               ],
               size=14, color=DARK)

    rect(sl, 6.85, 3.55, 6.05, 1.6, RGBColor(0xFD, 0xF2, 0xE9))
    bullet_box(sl, 7.05, 3.65, 5.7, 1.45,
               title="Excluded Dataset",
               items=[
                   "Merck/SMART: excluded per challenge rules",
               ],
               size=15, color=RED, title_color=RED)

    rect(sl, 6.85, 5.28, 6.05, 1.87, RGBColor(0xEA, 0xF4, 0xFF))
    bullet_box(sl, 7.05, 5.38, 5.7, 1.67,
               title="Why Not SENTRY?",
               items=[
                   "SENTRY also available via Vivli",
                   "ATLAS chosen: richer demographic metadata",
                   "Stronger Gram-negative coverage",
                   "Consistent MIC panel methodology",
               ],
               size=14, color=DARK)


def slide_06_data_quality(prs):
    """Data quality - censoring."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Data Quality Challenge: Left-Censoring",
               "80-90% of observations hit the panel floor - this is the most critical methodological issue")

    rect(sl, 0.4, 1.6, 5.8, 5.55, WHITE)
    bullet_box(sl, 0.6, 1.72, 5.5, 1.5,
               title="What is Left-Censoring?",
               items=[
                   "MIC panels have a fixed lowest concentration: 0.06 mg/L",
                   "Any isolate at or below this floor is reported as <=0.06",
                   "Most isolates are highly susceptible: ~80-90% hit the floor",
               ],
               size=14, color=DARK)

    rect(sl, 0.6, 3.2, 5.5, 0.06, TEAL)

    bullet_box(sl, 0.6, 3.35, 5.5, 3.55,
               title="How We Handle It",
               items=[
                   "<=0.06  -> 0.03 mg/L (half the floor)",
                   ">32     -> 32 mg/L (panel ceiling)",
                   "",
                   "is_censored flag added as model feature",
                   "pct_censored_year controls for panel shifts",
                   "",
                   "Key artifact: censoring dropped ~85% to ~25%",
                   "  during 2013-2017 then returned to ~88%",
                   "  Panel dilution range change - NOT biology",
                   "",
                   "3x sample weight on resistant isolates (MIC>=8)",
                   "  prevents model predicting floor for all obs.",
               ],
               size=13, color=DARK)

    # K. pneumoniae MIC90 trend chart - illustrates panel ceiling + censoring
    img_kp = PROJECT_ROOT / "reports" / "eda" / "mic90_trend_kpneumoniae_meropenem.png"
    if img_kp.exists():
        sl.shapes.add_picture(str(img_kp), Inches(6.5), Inches(1.55), Inches(6.6), Inches(5.6))
    else:
        rect(sl, 6.5, 1.55, 6.6, 5.6, WHITE)
        textbox(sl, 6.7, 4.0, 6.2, 0.5,
                "mic90_trend_kpneumoniae_meropenem.png", size=12, color=GRAY)


def slide_07_model(prs):
    """Modeling approach - interpretable baseline vs advanced model."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Model Development",
               "Interpretable baseline vs advanced model - explicit complexity/accuracy/interpretability tradeoff")

    # Left column - train/test split + target
    rect(sl, 0.4, 1.6, 6.0, 5.55, WHITE)
    bullet_box(sl, 0.6, 1.72, 5.6, 2.0,
               title="Train / Test Split (Time-Ordered)",
               items=[
                   "Train: 2004-2018  |  Test: 2019-2022",
                   "K. pneumoniae: 62,891 train / 26,681 test (70% / 30%)",
                   "A. baumannii:  24,003 train / 13,537 test (64% / 36%)",
                   "",
                   "NO random shuffling - would leak future data into training",
                   "and produce overly optimistic performance estimates.",
               ],
               size=13, color=DARK)

    rect(sl, 0.6, 3.72, 5.6, 0.04, TEAL)

    bullet_box(sl, 0.6, 3.85, 5.6, 1.55,
               title="Target Variable",
               items=[
                   "log2(MIC) - continuous regression",
                   "Log2 scale: 1 unit = 1 dilution step = minimum detectable difference",
                   "Evaluated: RMSE, MAE, R2 (full test) + RMSE/MAE on resistant subset",
               ],
               size=13, color=DARK)

    rect(sl, 0.6, 5.5, 5.6, 0.04, TEAL)

    bullet_box(sl, 0.6, 5.62, 5.6, 1.38,
               title="Why Not Linear Regression?",
               items=[
                   "Gene x country interactions are strongly non-linear",
                   "  (NDM in India vs NDM in Greece differ by 3+ log2 units)",
                   "Bimodal MIC target violates normality assumptions",
                   "Tree-based models are standard in AMR surveillance literature",
               ],
               size=12, color=GRAY, title_color=GRAY)

    # Right column - Model 1: RF (interpretable baseline)
    rect(sl, 6.8, 1.6, 6.1, 2.6, NAVY)
    bullet_box(sl, 7.0, 1.72, 5.7, 2.35,
               title="(1) Random Forest - More Interpretable Baseline",
               title_color=WHITE,
               items=[
                   "scikit-learn, n_estimators=200",
                   "Direct feature importances - no SHAP post-hoc needed",
                   "Each tree is a decision path reviewers can inspect",
                   "3x weight on resistant isolates (MIC >= 8 mg/L)",
                   "Better overall RMSE; weaker on resistant-subset precision",
               ],
               size=13, color=RGBColor(0xDD, 0xEE, 0xFF))

    # Model 2: XGBoost (advanced)
    rect(sl, 6.8, 4.35, 6.1, 2.75, RGBColor(0x0D, 0x3B, 0x6E))
    bullet_box(sl, 7.0, 4.47, 5.7, 2.5,
               title="(2) XGBoost + Optuna - Advanced Model",
               title_color=WHITE,
               items=[
                   "60 Optuna trials, internal validation: last 3 training years",
                   "Captures gene x country interaction effects",
                   "Tradeoff: +32% accuracy on resistant subset (K. pneu)",
                   "  vs interpretability cost - requires SHAP for explanation",
                   "Independent model per species",
               ],
               size=13, color=RGBColor(0xBB, 0xDD, 0xFF))

    rect(sl, 6.8, 7.18, 6.1, 0.22, TEAL)
    textbox(sl, 6.9, 7.21, 5.9, 0.19,
            "Tradeoff: XGBoost wins on clinically relevant resistant-subset RMSE; RF is more directly interpretable",
            size=10, color=WHITE, italic=True)


def slide_08_features(prs):
    """Feature set."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Feature Set",
               "Identical construction for both species - country dummy columns differ by coverage")

    features = [
        ("year",               "Continuous",   "Primary MIC creep driver"),
        ("gender_male",        "Binary",        "1 = Male; sex surveillance covariate"),
        ("age_paediatric",     "Binary",        "Age 0-17; adults (18-60) are reference"),
        ("age_elderly",        "Binary",        "Age 61+; adults are reference"),
        ("military_proxy",     "Binary",        "Wound/abscess + male + 18-60; combat-related proxy"),
        ("spec_*",             "OHE (5 levels)","Specimen: wound, blood, respiratory, urine, peritoneal"),
        ("ctry_*",             "OHE (65 cols)", "Country of isolation; Argentina = reference (drop_first)"),
        ("KPC_pos ... GES_pos","Binary x6",     "Carbapenemase genes: KPC, NDM, OXA, VIM, IMP, GES"),
        ("is_censored",        "Binary",        "DATA ARTIFACT - panel floor flag; not biological"),
        ("pct_censored_year",  "Float",         "Year-level censoring rate; controls panel methodology shifts"),
    ]

    rect(sl, 0.4, 1.55, 12.5, 0.42, NAVY)
    for i, (label, head) in enumerate([(0.6, "Feature"), (4.5, "Type"), (6.4, "Rationale")]):
        textbox(sl, label, 1.6, 3.5, 0.35, head, size=14, bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)

    for row, (feat, ftype, note) in enumerate(features):
        bg = WHITE if row % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF8)
        y = 1.97 + row * 0.48
        rect(sl, 0.4, y, 12.5, 0.48, bg)
        feat_color = RED if "censored" in feat else DARK
        textbox(sl, 0.55, y + 0.05, 3.7, 0.38, feat,
                size=13, bold=False, color=feat_color)
        textbox(sl, 4.45, y + 0.05, 1.75, 0.38, ftype,
                size=13, color=TEAL)
        textbox(sl, 6.35, y + 0.05, 6.3, 0.38, note,
                size=13, color=GRAY)

    textbox(sl, 0.4, 6.88, 12.5, 0.3,
            "is_censored and pct_censored_year are data-structure variables, not biological predictors - flagged explicitly in all SHAP plots",
            size=11, color=RED, italic=True)


def slide_09_results(prs):
    """Unified results table - both species."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Model Results: Both Species",
               "XGBoost vs Random Forest baseline - log2 MIC units - test period 2019-2022")

    # Shared column layout
    col_x = [0.45, 2.55, 4.35, 6.15, 7.85, 9.55, 11.15]
    col_w = [2.0,  1.7,  1.7,  1.6,  1.6,  1.5,  2.0]
    headers = ["Species / Model", "RMSE (all)", "MAE (all)", "RMSE (R)", "MAE (R)", "N resistant", "Improvement"]

    def draw_table(y_start, species_label, color,
                   rf_row, xgb_row, improvement_note):
        # Species header bar
        rect(sl, 0.4, y_start, 12.5, 0.38, color)
        textbox(sl, 0.55, y_start + 0.06, 5.0, 0.28,
                species_label, size=14, bold=True, color=WHITE, italic=True)

        for ri, (model_label, vals, row_color) in enumerate([
            ("RF baseline",   rf_row,  RGBColor(0xF4, 0xF6, 0xF9)),
            ("XGBoost tuned", xgb_row, WHITE),
        ]):
            y = y_start + 0.38 + ri * 0.48
            rect(sl, 0.4, y, 12.5, 0.48, row_color)
            textbox(sl, col_x[0], y + 0.1, col_w[0], 0.3,
                    model_label, size=13, bold=(ri == 1), color=DARK)
            for vi, (val, cx, cw) in enumerate(zip(vals, col_x[1:], col_w[1:])):
                is_best = (ri == 1) and vi in (1, 2)  # RMSE(R) and MAE(R) cols
                textbox(sl, cx, y + 0.1, cw, 0.3, val, size=13,
                        bold=is_best,
                        color=TEAL if is_best else DARK,
                        align=PP_ALIGN.CENTER)

        # Improvement note
        y_note = y_start + 0.38 + 2 * 0.48
        rect(sl, 0.4, y_note, 12.5, 0.32, RGBColor(0xE8, 0xF5, 0xE9))
        textbox(sl, 0.55, y_note + 0.06, 12.0, 0.22,
                improvement_note, size=11, color=RGBColor(0x1B, 0x6B, 0x3A), italic=True)

    # Column header row
    rect(sl, 0.4, 1.58, 12.5, 0.38, DARK)
    for hdr, cx, cw in zip(headers, col_x, col_w):
        textbox(sl, cx, 1.63, cw, 0.3, hdr, size=12, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER if cx > 1 else PP_ALIGN.LEFT)

    # K. pneumoniae block
    draw_table(
        y_start=1.96, species_label="Klebsiella pneumoniae + Meropenem",
        color=NAVY,
        rf_row =["1.558", "1.109", "2.869", "2.180", "4,305", ""],
        xgb_row=["1.758", "1.002", "1.960", "1.127", "4,305", "RMSE(R) -32%"],
        improvement_note="XGBoost reduces resistant-subset RMSE by 32% vs RF. "
                         "Higher overall RMSE expected: model correctly deprioritises the censoring floor.",
    )

    # A. baumannii block
    draw_table(
        y_start=4.08, species_label="Acinetobacter baumannii + Meropenem",
        color=RED,
        rf_row =["1.338", "0.789", "0.983", "0.510", "9,415", ""],
        xgb_row=["1.379", "0.707", "0.748", "0.270", "9,415", "RMSE(R) -24%"],
        improvement_note="XGBoost reduces resistant-subset RMSE by 24% vs RF. "
                         "Lower absolute RMSE than K. pneu because resistance is near-universal - less variance to explain.",
    )

    textbox(sl, 0.4, 6.68, 12.5, 0.52,
            "Metrics in log2 MIC units. RMSE (R) and MAE (R) = resistant subset (MIC >= 8 mg/L) - the clinically relevant measure. "
            "R2 on full test set is low for both species due to bimodal MIC distribution (80% of isolates at censoring floor); "
            "full metrics including R2 reported in supplementary notebook.",
            size=10, color=GRAY, italic=True)


def slide_10_shap(prs):
    """SHAP feature importance."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "SHAP Feature Importance: K. pneumoniae",
               "SHapley Additive exPlanations - how much each feature shifts the predicted log2(MIC), averaged across all test isolates")

    shap_data = [
        (1,  "KPC_pos",          0.943, "KPC carbapenemase - dominant mechanism 2004-2018", NAVY),
        (2,  "is_censored",      0.543, "DATA ARTIFACT - panel floor; not biological", RED),
        (3,  "OXA_pos",          0.504, "OXA-48/OXA-232 - dominant in Europe, Middle East", TEAL),
        (4,  "NDM_pos",          0.438, "NDM - bypasses avibactam; underrepresented in train", AMBER),
        (5,  "pct_censored_year",0.331, "Surveillance methodology control - not biological", RED),
        (6,  "ctry_China",       0.304, "High-resistance country effect", NAVY),
        (7,  "ctry_India",       0.160, "NDM origin country", NAVY),
        (8,  "ctry_Greece",      0.115, "OXA-48 hotspot", NAVY),
        (9,  "ctry_Turkey",      0.093, "OXA-48 hotspot", NAVY),
        (10, "year",             0.071, "Temporal MIC creep signal", TEAL),
    ]

    max_shap = 0.943
    bar_start = 4.3
    bar_max_w = 7.0

    rect(sl, 0.4, 1.55, 12.5, 0.42, NAVY)
    for label, x in [("Rank", 0.5), ("Feature", 1.1), ("Mean |SHAP|", 2.85), ("", bar_start)]:
        textbox(sl, x, 1.6, 2.5, 0.35, label, size=13, bold=True, color=WHITE)

    for i, (rank, feat, val, note, color) in enumerate(shap_data):
        y = 1.97 + i * 0.47
        bg = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF8)
        rect(sl, 0.4, y, 12.5, 0.47, bg)
        textbox(sl, 0.5, y + 0.08, 0.5, 0.3, str(rank), size=13, color=GRAY,
                align=PP_ALIGN.CENTER)
        feat_color = RED if "censored" in feat else DARK
        textbox(sl, 1.05, y + 0.08, 1.65, 0.3, feat, size=12, color=feat_color, bold=("censored" in feat))
        textbox(sl, 2.8, y + 0.08, 1.35, 0.3, f"{val:.3f}", size=13, color=color, bold=True,
                align=PP_ALIGN.CENTER)
        bar_w = (val / max_shap) * bar_max_w * 0.97
        rect(sl, bar_start, y + 0.1, bar_w, 0.26, color)
        textbox(sl, bar_start + bar_w + 0.05, y + 0.08, 12.5 - bar_start - bar_w - 0.4, 0.3,
                note, size=10, color=GRAY)

    textbox(sl, 0.4, 6.78, 12.5, 0.38,
            "KPC ranks above NDM because training period (2004-2018) is KPC-dominated. "
            "NDM only became prevalent after 2018 - this is a data recency limitation, not a model error. "
            "For A. baumannii, OXA-type genes are expected to rank highest.",
            size=11, color=GRAY, italic=True)


def slide_09b_comparison(prs):
    """Cross-species comparison and interpretation."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Cross-Species Comparison and Interpretation",
               "Two distinct resistance landscapes - same model architecture, very different epidemiology")

    # MIC90 comparison table
    rect(sl, 0.4, 1.58, 5.7, 5.55, WHITE)
    textbox(sl, 0.6, 1.65, 5.3, 0.38,
            "MIC90 Test Period (2019-2022)", size=15, bold=True, color=DARK)

    tbl_headers = ["Year", "Kp Actual", "Kp XGB", "Ab Actual", "Ab XGB"]
    tbl_cx = [0.55, 1.5, 2.6, 3.7, 4.75]
    tbl_cw = [0.85, 1.0, 1.0, 1.0, 1.0]
    tbl_rows = [
        ["2019", "32 mg/L", "~30", "32 mg/L", "~32"],
        ["2020", "32 mg/L", "~30", "32 mg/L", "~31"],
        ["2021", "32 mg/L", "~32", "32 mg/L", "~32"],
        ["2022", "32 mg/L", "~30", "32 mg/L", "~31"],
    ]
    rect(sl, 0.4, 2.05, 5.7, 0.38, DARK)
    for hdr, cx, cw in zip(tbl_headers, tbl_cx, tbl_cw):
        textbox(sl, cx, 2.1, cw, 0.28, hdr, size=12, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(tbl_rows):
        bg = RGBColor(0xF4, 0xF6, 0xF9) if ri % 2 == 0 else WHITE
        y = 2.43 + ri * 0.4
        rect(sl, 0.4, y, 5.7, 0.4, bg)
        for j, (cell, cx, cw) in enumerate(zip(row, tbl_cx, tbl_cw)):
            col = NAVY if j in (1, 2) else (RED if j in (3, 4) else DARK)
            textbox(sl, cx, y + 0.08, cw, 0.28, cell, size=12,
                    color=col, align=PP_ALIGN.CENTER)

    # Resistance rate comparison
    rect(sl, 0.4, 4.07, 5.7, 0.38, RGBColor(0x2A, 0x2A, 0x2A))
    textbox(sl, 0.55, 4.12, 5.4, 0.28,
            "Resistance rate in test period", size=12, bold=True, color=WHITE)
    resist_rows = [
        ("K. pneumoniae", "~19% (2019)", "~20% (2022)", NAVY),
        ("A. baumannii",  "~70% (2019)", "~67% (2022)", RED),
    ]
    for ri, (sp, v19, v22, col) in enumerate(resist_rows):
        y = 4.45 + ri * 0.42
        rect(sl, 0.4, y, 5.7, 0.42, RGBColor(0xF0, 0xF4, 0xF8) if ri % 2 == 0 else WHITE)
        textbox(sl, 0.55, y + 0.09, 2.1, 0.28, sp, size=12, bold=True, color=col, italic=True)
        textbox(sl, 2.65, y + 0.09, 1.2, 0.28, v19, size=12, color=DARK, align=PP_ALIGN.CENTER)
        textbox(sl, 3.85, y + 0.09, 1.2, 0.28, v22, size=12, color=DARK, align=PP_ALIGN.CENTER)

    textbox(sl, 0.55, 5.35, 5.4, 0.35,
            "Both at or near panel ceiling (32 mg/L) - true values may be higher",
            size=11, color=AMBER, italic=True)

    # Interpretation panels
    interp = [
        (NAVY, "K. pneumoniae - Creeping resistance",
         "MIC90 crossed EUCAST R threshold ~2017. "
         "Resistance tripled 2007-2024 (5% to 20%). "
         "KPC dominated training period; NDM rising fast post-2018. "
         "Model correctly tracks upward trend but undershoots ceiling."),
        (RED, "A. baumannii - Established crisis",
         "MIC90 at ceiling since 2005 - resistance was already entrenched "
         "before modern carbapenem surveillance. "
         "67-71% resistant throughout test period. "
         "OXA dominates. Year signal is rank #2 - temporal creep is real "
         "but superimposed on a high baseline."),
        (TEAL, "Key cross-species insight",
         "Lower absolute RMSE for A. baumannii (0.748 vs 1.960 on resistant subset) "
         "reflects a simpler prediction task - near-universal resistance means "
         "less variance. The harder problem is K. pneumoniae where the model "
         "must distinguish susceptible from resistant isolates."),
    ]
    for i, (col, title, body) in enumerate(interp):
        y = 1.58 + i * 1.82
        rect(sl, 6.55, y, 6.4, 1.72, col)
        textbox(sl, 6.75, y + 0.1, 6.0, 0.38,
                title, size=14, bold=True, color=WHITE)
        textbox(sl, 6.75, y + 0.5, 6.0, 1.15,
                body, size=12, color=RGBColor(0xEE, 0xEE, 0xEE))


def slide_10b_shap_ab(prs):
    """SHAP feature importance - A. baumannii."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "SHAP Feature Importance: A. baumannii",
               "Year ranks #2 - stronger temporal creep signal than K. pneumoniae (where year was #10)")

    shap_data = [
        (1,  "is_censored",        1.783, "DATA ARTIFACT - panel floor; not biological", RED),
        (2,  "year",               0.106, "Temporal MIC creep signal - stronger than K. pneu", TEAL),
        (3,  "pct_censored_year",  0.090, "Surveillance methodology control - not biological", RED),
        (4,  "age_paediatric",     0.063, "Paediatric isolates show distinct MIC distribution", AMBER),
        (5,  "OXA_pos",            0.056, "OXA carbapenemase - dominant mechanism in A. baumannii", NAVY),
        (6,  "ctry_France",        0.055, "Country effect: France", NAVY),
        (7,  "ctry_Italy",         0.051, "Country effect: Italy", NAVY),
        (8,  "ctry_Germany",       0.051, "Country effect: Germany", NAVY),
        (9,  "spec_wound",         0.046, "Wound specimen - military/trauma proxy", TEAL),
        (10, "spec_respiratory",   0.045, "Respiratory specimen - VAP association", TEAL),
    ]

    max_shap = 1.783
    bar_start = 4.3
    bar_max_w = 7.0

    rect(sl, 0.4, 1.55, 12.5, 0.42, RGBColor(0x6B, 0x0A, 0x0A))
    for label, x in [("Rank", 0.5), ("Feature", 1.1), ("Mean |SHAP|", 2.85), ("", bar_start)]:
        textbox(sl, x, 1.6, 2.5, 0.35, label, size=13, bold=True, color=WHITE)

    for i, (rank, feat, val, note, color) in enumerate(shap_data):
        y = 1.97 + i * 0.47
        bg = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF8)
        rect(sl, 0.4, y, 12.5, 0.47, bg)
        textbox(sl, 0.5, y + 0.08, 0.5, 0.3, str(rank), size=13, color=GRAY,
                align=PP_ALIGN.CENTER)
        feat_color = RED if "censored" in feat else DARK
        textbox(sl, 1.05, y + 0.08, 1.65, 0.3, feat, size=12,
                color=feat_color, bold=("censored" in feat))
        textbox(sl, 2.8, y + 0.08, 1.35, 0.3, f"{val:.3f}", size=13,
                color=color, bold=True, align=PP_ALIGN.CENTER)
        bar_w = (val / max_shap) * bar_max_w * 0.97
        rect(sl, bar_start, y + 0.1, bar_w if bar_w > 0.05 else 0.05, 0.26, color)
        textbox(sl, bar_start + bar_w + 0.1, y + 0.08,
                12.5 - bar_start - bar_w - 0.4, 0.3, note, size=10, color=GRAY)

    textbox(sl, 0.4, 6.78, 12.5, 0.38,
            "No KPC/NDM in A. baumannii - OXA-type carbapenemases (OXA-23, OXA-40, OXA-58) dominate throughout. "
            "Year at rank #2 confirms genuine temporal creep independent of gene prevalence. "
            "Croatia, Latvia, Lithuania, Ukraine are top resistance hotspots (94-98%).",
            size=11, color=GRAY, italic=True)


def _two_chart_slide(prs, title, subtitle, img_left, img_right, cap_left, cap_right,
                     col_left=None, col_right=None):
    """Reusable side-by-side chart slide."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, title, subtitle)
    col_left  = col_left  or NAVY
    col_right = col_right or RED
    EDA   = PROJECT_ROOT / "reports" / "eda"
    MODEL = PROJECT_ROOT / "reports" / "model"
    for img_path, x, cap, col in [
        (img_left,  0.3,  cap_left,  col_left),
        (img_right, 6.9,  cap_right, col_right),
    ]:
        # auto-detect subdirectory from filename pattern
        if any(img_path.startswith(p) for p in ("gene_", "mic90_trend", "specimen_")):
            p = EDA / img_path
        else:
            p = MODEL / img_path
        w = 6.2 if x < 5 else 6.1
        if p.exists():
            sl.shapes.add_picture(str(p), Inches(x), Inches(1.5), Inches(w), Inches(4.85))
        else:
            rect(sl, x, 1.5, w, 4.85, WHITE)
            textbox(sl, x + 0.2, 3.8, w - 0.4, 0.4, img_path, size=11, color=GRAY)
        rect(sl, x, 6.48, w, 0.72, col)
        textbox(sl, x + 0.2, 6.52, w - 0.4, 0.65, cap, size=12, color=WHITE)
    return sl


def slide_11a_gene_prevalence(prs):
    """Gene prevalence over time - both species."""
    _two_chart_slide(
        prs,
        title="Carbapenemase Gene Prevalence Over Time",
        subtitle="Dominant resistance mechanisms differ sharply between the two species",
        img_left="gene_prevalence_over_time_kpneumoniae.png",
        img_right="gene_prevalence_over_time_abaumannii.png",
        cap_left="K. pneumoniae: KPC dominated 2004-2016, then NDM + OXA rose sharply post-2017",
        cap_right="A. baumannii: OXA dominates throughout - no KPC; NDM secondary",
        col_left=NAVY,
        col_right=RED,
    )


def slide_11b_specimen_source(prs):
    """Specimen source MIC90 - both species."""
    _two_chart_slide(
        prs,
        title="MIC90 by Specimen Source",
        subtitle="Both species: all 6 specimen types shown. K. pneumoniae shows variation; A. baumannii saturates at 32 mg/L across all sources",
        img_left="specimen_source_mic90_kpneumoniae.png",
        img_right="specimen_source_mic90_abaumannii.png",
        cap_left="K. pneumoniae: respiratory/blood/peritoneal at MIC90=8 mg/L; urine lowest (1 mg/L)",
        cap_right="A. baumannii: all sources at ceiling (32 mg/L); bars colored by %R - respiratory highest (70%R)",
        col_left=NAVY,
        col_right=RED,
    )


def slide_11c_rmse(prs):
    """RMSE by year - both species."""
    _two_chart_slide(
        prs,
        title="Model Performance: RMSE by Year (Test 2019-2022)",
        subtitle="XGBoost outperforms Random Forest baseline for resistant isolates in both species",
        img_left="rmse_by_year_kpneumoniae.png",
        img_right="rmse_by_year_abaumannii.png",
        cap_left="K. pneumoniae: XGBoost RMSE(R) 1.960 vs RF 2.869 - 32% improvement",
        cap_right="A. baumannii: XGBoost RMSE(R) 0.748 vs RF 0.983 - 24% improvement",
        col_left=NAVY,
        col_right=RED,
    )


def slide_11d_residuals(prs):
    """Residuals by year - both species."""
    _two_chart_slide(
        prs,
        title="Residuals by Year (XGBoost, Test Set)",
        subtitle="Residual pattern reveals where the model under- or over-predicts",
        img_left="residuals_by_year_kpneumoniae.png",
        img_right="residuals_by_year_abaumannii.png",
        cap_left="K. pneumoniae: positive residuals post-2018 - model undershoots the ceiling",
        cap_right="A. baumannii: residuals near zero - MIC90 already at ceiling throughout",
        col_left=NAVY,
        col_right=RED,
    )


def slide_11e_shap_charts(prs):
    """SHAP beeswarm plots - both species."""
    _two_chart_slide(
        prs,
        title="SHAP Feature Importance: Beeswarm Plots",
        subtitle="Each dot = one isolate from the test set. Position on x-axis = how much that feature pushed the predicted MIC up (right) or down (left). Colour = feature value: red = high, blue = low.",
        img_left="shap_beeswarm_kpneumoniae.png",
        img_right="shap_beeswarm_abaumannii.png",
        cap_left="K. pneumoniae: KPC, OXA, NDM are top biological drivers; year at rank #10",
        cap_right="A. baumannii: OXA dominates; year at rank #2 - stronger temporal creep signal",
        col_left=NAVY,
        col_right=RED,
    )


def slide_11_limitations(prs):
    """Known limitations."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Known Limitations",
               "Transparent reporting of methodological constraints")

    limitations = [
        (RED,   "Panel Ceiling (32 mg/L)",
                "MIC90 saturates post-2018. True values may be 64-128+ mg/L. "
                "Trend analysis is valid for 2004-2018 training period. "
                "Linear slope underestimates true escalation rate."),
        (AMBER, "Bimodal Target",
                "~75% of observations cluster at the censoring floor. "
                "RMSE dominated by floor spike without sample weighting. "
                "3x resistant upweight partially mitigates - full solution requires tobit regression."),
        (NAVY,  "NDM Underrepresented in Train",
                "NDM rose to dominance post-2018 - after the training cutoff. "
                "The model is less sensitive to the fastest-rising resistance mechanism. "
                "Stated explicitly; will improve with post-2020 retraining."),
        (TEAL,  "Military Proxy is Broad",
                "wound + male + 18-60 captures 2,813 isolates (2.8%) of which many are not combat-related. "
                "Stated as approximate; acceptable for proof-of-concept scope."),
        (GRAY,  "Paediatric MIC90 Artifact",
                "Flat at 0.12 mg/L (2006-2016), then jumps to 32 in 2019. "
                "Driven by censoring floor (91.7% at floor pre-2017) combined with small n. "
                "Excluded from primary MIC creep trend analysis."),
    ]

    for i, (color, title, body) in enumerate(limitations):
        y = 1.6 + i * 1.12
        rect(sl, 0.4, y, 0.18, 0.9, color)
        rect(sl, 0.58, y, 12.32, 0.9, WHITE)
        textbox(sl, 0.75, y + 0.05, 3.2, 0.38, title,
                size=15, bold=True, color=color)
        textbox(sl, 0.75, y + 0.45, 11.9, 0.38, body,
                size=13, color=DARK)


def slide_12_api(prs):
    """API + Dashboard."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "API and Dashboard",
               "FastAPI inference service + Next.js frontend - examples below show K. pneumoniae (default species)")

    rect(sl, 0.4, 1.6, 6.0, 5.55, WHITE)
    bullet_box(sl, 0.6, 1.72, 5.6, 0.45,
               title="FastAPI Endpoints",
               items=[], size=16)

    endpoints = [
        ("GET /health",                 "Service status, loaded models"),
        ("GET /methodology",            "Methodology page (HTML)"),
        ("GET /api/trend/mic90",        "MIC90 by year: actual + predicted + forecast"),
        ("GET /api/country-stats",      "Resistance rate and MIC90 by country"),
        ("GET /api/features/importance","Top 20 SHAP features"),
        ("POST /api/predict",           "Single-isolate MIC prediction"),
        ("GET /api/countries",          "Countries known to the model"),
    ]
    for i, (ep, desc) in enumerate(endpoints):
        y = 2.25 + i * 0.7
        rect(sl, 0.6, y, 5.6, 0.65, RGBColor(0xF4, 0xF6, 0xF9) if i % 2 == 0 else WHITE)
        textbox(sl, 0.75, y + 0.04, 2.8, 0.28, ep, size=11, color=NAVY, bold=True)
        textbox(sl, 0.75, y + 0.32, 5.2, 0.28, desc, size=12, color=GRAY)

    rect(sl, 6.8, 1.6, 6.1, 2.55, NAVY)
    bullet_box(sl, 7.0, 1.72, 5.7, 2.3,
               title="All Endpoints Accept ?species=",
               title_color=WHITE,
               items=[
                   "kpneumoniae  -  K. pneumoniae model",
                   "abaumannii   -  A. baumannii model",
                   "",
                   "API gracefully returns 503 if",
                   "species model not yet loaded",
               ],
               size=14, color=RGBColor(0xBB, 0xDD, 0xFF))

    rect(sl, 6.8, 4.3, 6.1, 2.85, RGBColor(0x12, 0x4A, 0x42))
    bullet_box(sl, 7.0, 4.42, 5.7, 2.6,
               title="Infrastructure (Free Tier)",
               title_color=WHITE,
               items=[
                   "Code:         GitHub (public, open-source)",
                   "Model files:  Hugging Face Hub",
                   "Aggregated DB: Supabase (PostgreSQL)",
                   "API:          Render.com",
                   "Frontend:     Vercel (Next.js, auto-deploy)",
                   "",
                   "Raw data:     Local only - never uploaded",
               ],
               size=13, color=RGBColor(0xAA, 0xFF, 0xDD))


def slide_14_dashboard(prs):
    """Web dashboard showcase."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT_BG)
    header_bar(sl, "Live Web Dashboard",
               "https://mic-creep-predict.vercel.app  -  publicly accessible, no login required")

    # Left: 4 feature cards in 2x2 grid
    cards = [
        (NAVY,  "MIC Trend",
                "Year-by-year MIC90 + resistance rate.\n"
                "Actual vs XGBoost predictions, 2004-2022."),
        (TEAL,  "Countries",
                "MIC90 and resistance rate across\n81 countries, species-switchable."),
        (RGBColor(0x12, 0x5A, 0x50), "Feature Importance",
                "Top 20 SHAP features. Data artifact\nfeatures explicitly flagged."),
        (RED,   "Methodology",
                "Full write-up: data, model, features,\ncensoring, limitations, GitHub link."),
    ]
    for i, (color, title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 3.05
        y = 1.55 + row * 2.72
        w = 2.85
        h = 2.52
        rect(sl, x, y, w, h, color)
        textbox(sl, x + 0.15, y + 0.12, w - 0.25, 0.42,
                title, size=16, bold=True, color=WHITE)
        textbox(sl, x + 0.15, y + 0.58, w - 0.25, 1.75,
                body, size=12, color=RGBColor(0xEE, 0xEE, 0xEE))

    # Right: MIC Predictor screenshot
    predictor_img = PROJECT_ROOT / "MIC Predictor.png"
    if predictor_img.exists():
        sl.shapes.add_picture(str(predictor_img), Inches(6.5), Inches(1.5), Inches(6.55), Inches(5.28))
    else:
        rect(sl, 6.5, 1.5, 6.55, 5.28, WHITE)
        textbox(sl, 6.7, 3.8, 6.2, 0.4, "MIC Predictor.png", size=12, color=GRAY)

    # Caption under screenshot
    rect(sl, 6.5, 6.82, 6.55, 0.52, NAVY)
    textbox(sl, 6.65, 6.88, 6.3, 0.42,
            "Predict page: Ukraine, Female, Respiratory, NDM + OXA-48 + GES  ->  11.91 mg/L  RESISTANT",
            size=12, bold=True, color=WHITE)

    textbox(sl, 0.3, 7.18, 6.1, 0.28,
            "Deployed on Vercel - auto-deploys from GitHub on every push",
            size=10, color=GRAY, italic=True)


def slide_13_conclusions(prs):
    """Conclusions and Lessons Learned."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, DARK)
    rect(sl, 0, 0, 13.33, 0.12, TEAL)

    textbox(sl, 0.7, 0.2, 12, 0.55,
            "Conclusions and Lessons Learned", size=28, bold=True,
            color=WHITE, align=PP_ALIGN.LEFT)

    # Left column: Findings + Best model + Practical implications
    rect(sl, 0.4, 0.95, 6.1, 6.4, RGBColor(0x0D, 0x1B, 0x2A))

    bullet_box(sl, 0.6, 1.05, 5.7, 2.6,
               title="Main Findings",
               title_color=TEAL,
               items=[
                   "K. pneumoniae: measurable MIC creep confirmed",
                   "  +1.97 mg/L/yr MIC90 slope (R2=0.67, p<0.001)",
                   "  Resistance tripled: 5% (2007) -> 20% (2022)",
                   "  Model detects the trend before threshold breach",
                   "",
                   "A. baumannii: resistance already at crisis level",
                   "  MIC90 at panel ceiling (32 mg/L) since 2005",
                   "  70% resistant throughout test period 2019-2022",
               ],
               size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    bullet_box(sl, 0.6, 3.75, 5.7, 1.35,
               title="Best Model",
               title_color=TEAL,
               items=[
                   "XGBoost (tuned): -32% RMSE on resistant K. pneu isolates vs RF",
                   "RF provides more direct interpretability (feature importances)",
                   "  without SHAP - preferred when transparency is critical",
               ],
               size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    bullet_box(sl, 0.6, 5.2, 5.7, 1.0,
               title="Practical Implications",
               title_color=TEAL,
               items=[
                   "Enables early-warning surveillance before thresholds are breached",
                   "Live API + dashboard: publicly accessible, no login required",
                   "SHAP outputs require domain-expert validation before clinical use",
               ],
               size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Right column: Lessons learned + Future improvements
    rect(sl, 6.85, 0.95, 6.1, 6.4, RGBColor(0x0D, 0x1B, 0x2A))

    bullet_box(sl, 7.05, 1.05, 5.7, 3.15,
               title="Lessons Learned",
               title_color=AMBER,
               items=[
                   "RMSE on full dataset misleads for bimodal MIC data",
                   "  80% at censoring floor -> RMSE(resistant) is primary metric",
                   "",
                   "Time-aware splits are non-negotiable for surveillance data",
                   "  random shuffle produces artificially optimistic estimates",
                   "",
                   "Panel artifacts need explicit features in the model",
                   "  is_censored and pct_censored_year prevent confounding",
                   "",
                   "Gene x country interactions require tree-based models",
                   "  NDM in India vs NDM in Greece differ by 3+ log2 units",
               ],
               size=12, color=RGBColor(0xFF, 0xEE, 0xCC))

    bullet_box(sl, 7.05, 4.3, 5.7, 2.9,
               title="Future Improvements",
               title_color=AMBER,
               items=[
                   "Retrain with post-2022 data: NDM rising fast post-2018",
                   "Tobit regression: principled handling of censored values",
                   "Extend to all WHO priority pathogens and antibiotics",
                   "Automated retraining pipeline for ongoing surveillance",
                   "Integrate with Ukraine national AMR surveillance infrastructure",
               ],
               size=12, color=RGBColor(0xFF, 0xEE, 0xCC))


# =============================================================================
# Build
# =============================================================================

def main():
    prs = new_prs()
    # --- Motivation ---
    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_why_matters(prs)
    slide_04_pathogens(prs)
    # --- Data ---
    slide_05_data(prs)
    # --- EDA ---
    slide_11a_gene_prevalence(prs)
    slide_11b_specimen_source(prs)
    slide_06_data_quality(prs)
    # --- Methods ---
    slide_08_features(prs)
    slide_07_model(prs)
    # --- Results ---
    slide_09_results(prs)
    slide_11c_rmse(prs)
    slide_11d_residuals(prs)
    # --- Interpretation ---
    slide_09b_comparison(prs)
    slide_10_shap(prs)
    slide_10b_shap_ab(prs)
    slide_11e_shap_charts(prs)
    # --- Limitations + Conclusions ---
    slide_11_limitations(prs)
    slide_13_conclusions(prs)
    # --- Appendix ---
    slide_12_api(prs)
    slide_14_dashboard(prs)

    prs.save(OUT)
    print(f"Saved: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
